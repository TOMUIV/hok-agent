import sys, os, json, re, math
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

BASE = r"D:\资料库\文档\THU\上课！\AI实践基石\HOK"
PPT_DIR = os.path.join(BASE, "ppt")
OUTPUT = os.path.join(PPT_DIR, "答辩PPT.pptx")

BACK_DEFAULT = os.path.join(PPT_DIR, "back.jpg")
LR_BACK = os.path.join(PPT_DIR, "LR_BACK.png")
PROOF_BACK = os.path.join(PPT_DIR, "PROOF_BACK.png")
AGENT_PNG = os.path.join(PPT_DIR, "agent.png")
PROMPT_PNG = os.path.join(PPT_DIR, "prompt.png")
TIG_PNG = os.path.join(PPT_DIR, "TIG.png")
WIA_PNG = os.path.join(PPT_DIR, "WIA.png")
REFLEXION_PNG = os.path.join(PPT_DIR, "reflexion.png")
TRAJ_SCREENSHOT = os.path.join(BASE, "src", "screenshot_traj.png")

NAVY = RGBColor(0x1B, 0x3A, 0x5C)
ORANGE = RGBColor(0xF2, 0x99, 0x45)  # softer orange for dark & light themes
WARM = RGBColor(0xF8, 0xF6, 0xF1)
DARK = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xE0, 0xE0, 0xE0)
MGRAY = RGBColor(0x99, 0x99, 0x99)
DGRAY = RGBColor(0x66, 0x66, 0x66)
CARD_BG = WHITE
CARD_BD = LGRAY
MASK_COLOR = RGBColor(0x00, 0x00, 0x00)

SW = Inches(13.333)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

SECTIONS = ["问题", "方案", "实现", "评估"]
PROGRESS_MAP = {
    3: (0, 1, 0),
    4: (1, 3, 0), 5: (1, 3, 1), 6: (1, 3, 2),
    7: (2, 3, 0), 8: (2, 3, 1), 9: (2, 3, 2),
    10: (3, 5, 0), 11: (3, 5, 1), 12: (3, 5, 2), 13: (3, 5, 3), 14: (3, 5, 4),
}

def is_cjk(c):
    return ("\u4e00" <= c <= "\u9fff" or "\uff00" <= c <= "\uffef" or "\u3000" <= c <= "\u303f")

def set_bg(slide, color=WARM):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def set_bg_img(slide, img_path):
    if os.path.isfile(img_path):
        slide.shapes.add_picture(img_path, 0, 0, SW, SH)

def _set_alpha(shape, alpha_permille):
    sp = shape._element
    spPr = sp.find(qn('a:spPr'))
    if spPr is None: spPr = sp.find(qn('p:spPr'))
    if spPr is None: return
    sf = spPr.find(qn('a:solidFill'))
    if sf is None: return
    srgb = sf.find(qn('a:srgbClr'))
    if srgb is None: return
    al = srgb.find(qn('a:alpha'))
    if al is None:
        al = etree.SubElement(srgb, qn('a:alpha'))
    al.set('val', str(alpha_permille * 100))

def add_mask(slide, mask_type='black', alpha_permille=300):
    color = RGBColor(0, 0, 0) if mask_type == 'black' else RGBColor(0xFF, 0xFF, 0xFF)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    _set_alpha(s, alpha_permille)

def add_card(slide, l, t, w, h):
    """White card with thin border."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = CARD_BG
    s.line.color.rgb = CARD_BD; s.line.width = Pt(0.5)
    return s

def add_card_st(slide, l, t, w, h, alpha=700):
    """Semi-transparent card with thin border (for light/white mask)."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0xF8, 0xF6, 0xF1)
    _set_alpha(s, alpha)
    s.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC); s.line.width = Pt(0.5)
    return s

def add_card_dark(slide, l, t, w, h, alpha=750):
    """Semi-transparent dark card (for black mask)."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x32)
    _set_alpha(s, alpha)
    s.line.color.rgb = RGBColor(0x44, 0x44, 0x55); s.line.width = Pt(0.5)
    return s

def add_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def add_tb(slide, l, t, w, h, text, sz, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    lines = text.split("\n")
    for li, line in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(1)
        segs = re.split(r"([\u4e00-\u9fff\uff00-\uffef\u3000-\u303f]+|[A-Za-z0-9._()%+\-→←≤≥×÷/]+)", line)
        for seg in segs:
            if not seg: continue
            r = p.add_run(); r.text = seg; r.font.size = Pt(sz)
            r.font.color.rgb = color; r.font.bold = bold
            r.font.name = "SimSun" if any(is_cjk(c) for c in seg) else "Times New Roman"
    return tb

def add_header(slide, label):
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.50))
    h.fill.solid(); h.fill.fore_color.rgb = RGBColor(0, 0, 0); h.line.fill.background()
    _set_alpha(h, 550)
    add_tb(slide, Inches(0.4), Inches(0.05), Inches(12), Inches(0.42), label, 18, WHITE, bold=True)

def add_progress(slide, slide_idx):
    if slide_idx <= 2 or slide_idx >= 13:
        return
    info = PROGRESS_MAP.get(slide_idx)
    if not info:
        return
    sect_idx, total_s, page_idx = info
    progress = 1.0 if slide_idx == 14 else (sect_idx / 4 + (page_idx + 0.5) / total_s / 4)
    bar_l, bar_t = Inches(0.5), SH - Inches(0.42)
    bar_w, bar_h = Inches(8.0), Inches(0.16)
    add_rect(slide, bar_l, bar_t, bar_w, bar_h, RGBColor(0x44, 0x44, 0x55))
    fill_w = int(bar_w * progress)
    if fill_w > Inches(0.05):
        add_rect(slide, bar_l, bar_t, fill_w, bar_h, ORANGE)
    seg_w = bar_w / 4
    for i, name in enumerate(SECTIONS):
        x = bar_l + i * seg_w + seg_w / 2 - Inches(0.35)
        add_tb(slide, x, bar_t + Inches(0.20), Inches(0.7), Inches(0.18), name, 8,
               ORANGE if i == sect_idx else RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)

def add_img(slide, path, l, t, w=None, h=None):
    if not os.path.isfile(path):
        return None
    kw = {"left": l, "top": t}
    if w: kw["width"] = w
    if h: kw["height"] = h
    return slide.shapes.add_picture(path, **kw)

def add_table(slide, data, left, top, width, height, dark=False):
    rows_n, cols_n = len(data), len(data[0])
    tbl_shape = slide.shapes.add_table(rows_n, cols_n, left, top, width, height)
    tbl = tbl_shape.table
    col_w = int(width / cols_n)
    for ci in range(cols_n):
        tbl.columns[ci].width = col_w
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci)
            c.text = str(val)
            # Vertical center
            tcPr = c._tc.get_or_add_tcPr()
            tcPr.set('anchor', 'ctr')
            for p in c.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(14) if ri == 0 else Pt(12)
                    run.font.bold = ri == 0
                    run.font.name = "SimSun" if any(is_cjk(ch) for ch in str(val)) else "Times New Roman"
                    run.font.color.rgb = WHITE if ri == 0 else DARK
            if ri == 0:
                c.fill.solid(); c.fill.fore_color.rgb = NAVY if not dark else ORANGE
            elif ri % 2 == 1:
                c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xF0, 0xEE, 0xE8)
    return tbl_shape

def make_slide(bg_img_path=None, header_text="", slide_idx=0, mask_type=None):
    s = prs.slides.add_slide(BLANK)
    if bg_img_path and os.path.isfile(bg_img_path):
        set_bg_img(s, bg_img_path)
        if mask_type:
            alpha = 350 if mask_type == 'black' else 450
            add_mask(s, mask_type, alpha)
    else:
        set_bg(s, WARM)
    if header_text:
        add_header(s, header_text)
    add_progress(s, slide_idx)
    return s

# ══════════ S1: Title ══════════
def slide_title():
    s = make_slide(BACK_DEFAULT, "", 1, 'black')
    add_tb(s, Inches(1.0), Inches(1.2), Inches(11.3), Inches(1.0),
           "Honor of Kings LLM Agent", 42, WHITE, bold=True, align=PP_ALIGN.CENTER)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.0), Inches(2.25), Inches(9.3), Inches(0.55))
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0, 0, 0); r.line.fill.background()
    _set_alpha(r, 400)
    add_tb(s, Inches(2.0), Inches(2.28), Inches(9.3), Inches(0.5),
           "基于提示词工程与分层记忆系统的纯文本 LLM 游戏智能体", 22, ORANGE, align=PP_ALIGN.CENTER)
    info_y = Inches(4.5)
    add_tb(s, Inches(2.0), info_y, Inches(9.3), Inches(0.3),
           "清华大学  人工智能学院  AI实践基石", 16, RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
    add_tb(s, Inches(2.0), info_y + Inches(0.35), Inches(9.3), Inches(0.3),
           "杨坤鑫 · 张腾达", 16, RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
    add_tb(s, Inches(2.0), info_y + Inches(0.70), Inches(9.3), Inches(0.3),
           "2026年7月17日", 14, MGRAY, align=PP_ALIGN.CENTER)

# ══════════ S2: TOC ══════════
def slide_toc():
    s = make_slide(BACK_DEFAULT, "目  录", 2, 'black')
    cards = [
        ("一、问题", "声明性与过程性知识的鸿沟",
         "LLM 会解释不会做 vs RL 会做不会解释\n"
         "TiG / WiA-LLM / Reflexion 的局限\n"
         "→ 零训练成本弥合鸿沟"),
        ("二、方案", "四系统流水线总体设计",
         "SYS1: 局内决策 (WhatIf → @SKILL_CALL)\n"
         "SYS2: 逐事件分析 → SYS3: 全局复盘\n"
         "→ AUDIT: 经验审计 → DB 持久化"),
        ("三、实现", "三大子系统底层细节",
         "三层记忆系统 + 去重引擎\n"
         "Prompt 继承架构与数据模型\n"
         "→ Macro/Concrete 双层技能系统"),
        ("四、评估", "实验验证与未来方向",
         "轨迹回放浏览器运行演示\n"
         "相关工作对比与局限分析\n"
         "→ 3v3 多 Agent + 聊天室对抗"),
    ]
    cw, ch = Inches(5.4), Inches(2.8)
    sx, sy = Inches(0.7), Inches(1.1)
    gap_x, gap_y = Inches(6.5), Inches(3.3)
    for i, (t, d, dd) in enumerate(cards):
        x = sx if i % 2 == 0 else sx + gap_x
        y = sy if i < 2 else sy + gap_y
        add_card_dark(s, x, y, cw, ch)
        add_tb(s, x + Inches(0.3), y + Inches(0.25), cw - Inches(0.6), Inches(0.35), t, 20, ORANGE, bold=True)
        add_tb(s, x + Inches(0.3), y + Inches(0.75), cw - Inches(0.6), Inches(0.3), d, 15, RGBColor(0xBB, 0xBB, 0xBB))
        add_tb(s, x + Inches(0.3), y + Inches(1.25), cw - Inches(0.6), Inches(1.3), dd, 14, WHITE)

# ══════════ S3: Problem ══════════
def slide_problem():
    s = make_slide(LR_BACK, "一·问题 — 声明性与过程性知识的鸿沟", 3, 'white')

    refs = [
        (TIG_PNG, "TiG (Liao et al., 2025)",
         "将 RL 决策重写为语言建模任务\n"
         "方法：SFT + GRPO 两阶段训练\n"
         "训练成本：需 8×H100 GPU 数千步\n"
         "记忆：无显式记忆，每局独立\n"
         "可解释：<think> 推理 + 宏观动作选择"),
        (WIA_PNG, "WiA-LLM (Sui et al., 2026)",
         "训练 LLM 作为显式语言世界模型\n"
         "方法：SFT + GRPO 两阶段训练\n"
         "训练成本：需 8×H100 GPU\n"
         "记忆：无显式记忆\n"
         "可解释：What-If 自然语言预测"),
        (REFLEXION_PNG, "Reflexion (Shinn et al., 2023)",
         "语言反馈 + Episodic Buffer 机制\n"
         "方法：口头反思 + 情节记忆\n"
         "训练成本：零训练\n"
         "记忆：单层 Episodic Buffer\n"
         "可解释：反思文本 + 改进决策"),
    ]
    card_w = Inches(4.3)
    card_gap = Inches(0.3)
    total_w = 3 * card_w + 2 * card_gap
    start_x = (SW - total_w) / 2
    for i, (img_path, title, desc) in enumerate(refs):
        x = start_x + i * (card_w + card_gap)
        add_card_st(s, x, Inches(0.9), card_w, Inches(3.3))
        add_tb(s, x + Inches(0.15), Inches(1.0), card_w - Inches(0.3), Inches(0.25), title, 13, NAVY, bold=True)
        add_img(s, img_path, x + Inches(0.15), Inches(1.35), w=card_w - Inches(0.3))
        add_tb(s, x + Inches(0.15), Inches(2.6), card_w - Inches(0.3), Inches(1.5), desc, 11, DARK)

    add_card_st(s, Inches(0.3), Inches(4.5), Inches(12.7), Inches(1.0))
    add_tb(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.25),
           "核心问题", 15, NAVY, bold=True)
    add_tb(s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.5),
           "RL 智能体擅长过程性知识（会做不会解释），LLM 擅长声明性知识（会解释不会做）。"
           "两者之间存在根本性鸿沟。本项目的目标是：在不依赖任何模型训练的前提下，"
           "通过纯提示词工程 + 分层记忆系统弥合这一鸿沟。", 12, DARK)

    add_card_st(s, Inches(0.3), Inches(5.75), Inches(12.7), Inches(1.2))
    add_tb(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.25),
           "我们的方案", 15, NAVY, bold=True)
    add_tb(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.65),
           "综上，现有的工作要么依赖大规模 GPU 训练（TiG、WiA-LLM），要么仅支持单层记忆（Reflexion），"
           "难以在零训练成本下兼顾决策质量与跨对局学习能力。本项目的核心创新在于：以纯提示词工程为基础，"
           "引入分层记忆系统和赛后反思流水线，在完全不进行模型训练的前提下弥合声明性与过程性知识的鸿沟。", 11, DARK)

# ══════════ S4: Architecture ══════════
def slide_arch():
    s = make_slide(BACK_DEFAULT, "二·方案 — 四系统流水线总览", 4, 'black')
    add_img(s, AGENT_PNG, (SW - Inches(5.91)) / 2, Inches(0.9), w=Inches(5.91))
    labels = [
        ("SYS1: 局内决策", Inches(8.5), Inches(1.5), "每帧 LLM WhatIf → @SKILL_CALL"),
        ("SYS2: 事件分析", Inches(8.5), Inches(2.5), "赛后逐事件深度复盘 → BUFFER"),
        ("SYS3: 全局复盘", Inches(8.5), Inches(3.3), "赛后整局总结 → BUFFER"),
        ("AUDIT: 经验审计", Inches(8.5), Inches(4.1), "评分+1/-1/0 → DB merge+去重"),
    ]
    for title, x, y, desc in labels:
        add_card_dark(s, x, y, Inches(4.3), Inches(0.55))
        add_tb(s, x + Inches(0.1), y + Inches(0.03), Inches(4.3), Inches(0.22), title, 14, ORANGE, bold=True)
        add_tb(s, x + Inches(0.1), y + Inches(0.27), Inches(4.3), Inches(0.26), desc, 14, RGBColor(0xBB,0xBB,0xBB))

# ══════════ S5: SYS1 ══════════
def slide_sys1():
    s = make_slide(BACK_DEFAULT, "二·方案 — SYS1: 局内决策循环", 5, 'black')
    flow_items = [
        ("Game State", "protobuf per frame"),
        ("parse_state()", "→ text state + MACRO ACTIONS"),
        ("LLM WhatIf", "<think> 5段推理"),
        ("@SKILL_CALL", "<action> 高层指令"),
        ("Concrete Skill", "多帧循环执行"),
        ("6-tuple", "(btn, mx, mz, sx, sz, target)"),
        ("Gamecore", "执行并前进一帧"),
    ]
    flow_y = Inches(1.0)
    for i, (title, desc) in enumerate(flow_items):
        y = flow_y + i * Inches(0.75)
        add_card_dark(s, Inches(0.5), y, Inches(3.2), Inches(0.62))
        add_tb(s, Inches(0.65), y + Inches(0.05), Inches(2.9), Inches(0.24), title, 14, ORANGE, bold=True)
        if desc:
            add_tb(s, Inches(0.65), y + Inches(0.32), Inches(2.9), Inches(0.24), desc, 13, RGBColor(0xBB,0xBB,0xBB))
        if i < len(flow_items) - 1:
            add_tb(s, Inches(1.8), y + Inches(0.55), Inches(0.6), Inches(0.18), "↓", 14, MGRAY, align=PP_ALIGN.CENTER)

    rx = Inches(4.3)
    add_card_dark(s, rx, Inches(0.95), Inches(2.0), Inches(0.4))
    add_tb(s, rx + Inches(0.1), Inches(0.97), Inches(1.8), Inches(0.35), "关键设计", 16, ORANGE, bold=True)
    points = [
        ("WhatIf 双候选推理",
         "每帧评估2个候选动作，预测各自结果后决策。借鉴 WiA-LLM 的 lookahead 思想。"),
        ("DELTA 六维追踪",
         "HP / 金钱 / 塔血 / 小兵 / 装备 / 增益。SELF 与 ENEMY 完全对称。"),
        ("@SKILL_CALL 抽象层",
         "LLM 不直接接触 6-tuple。每次输出高层指令，executor 自动翻译执行。"),
        ("Memory 注入",
         "每帧 USER 消息含 TRENDS + DETAIL + DELTA。跨局经验由 memory.retrieve() 注入。"),
    ]
    for i, (title, desc) in enumerate(points):
        y = Inches(1.5) + i * Inches(1.05)
        add_card_dark(s, rx, y, Inches(8.5), Inches(0.9))
        add_tb(s, rx + Inches(0.15), y + Inches(0.08), Inches(8.2), Inches(0.22), title, 15, ORANGE, bold=True)
        add_tb(s, rx + Inches(0.15), y + Inches(0.38), Inches(8.2), Inches(0.45), desc, 14, WHITE)

# ══════════ S6: Reflection ══════════
def slide_reflection():
    s = make_slide(BACK_DEFAULT, "二·方案 — 赛后反思流水线", 6, 'black')
    pipeline = [
        ("轨迹 JSONL",
         "输入：逐帧 protobuf 游戏状态\n"
         "存储：JSON Lines 格式每步写入\n"
         "字段：system_prompt / user_msg\n"
         "      llm_reply / parsed_results\n"
         "用途：赛后回溯、事件检测、DEBUG"),
        ("Event Detection",
         "输入：轨迹 JSONL 全量数据\n"
         "检测规则：HP 从 >0 变 0 → kill/death\n"
         "           金币变化 ≥200 → gold_spike\n"
         "           塔 HP 归零 → tower_fall\n"
         "输出：事件列表 [{type, frame}]"),
        ("SYS2 (×N)",
         "输入：单事件 + BEFORE(100帧)\n"
         "       + AFTER(100帧)\n"
         "处理：LLM 分析因果 + 提炼教训\n"
         "输出：新 EPISODIC (Context+Lesson)\n"
         "      新 SEMANTIC (规则) → BUFFER"),
        ("SYS3 (×1)",
         "输入：整局全帧 DETAIL\n"
         "处理：LLM 全局宏观视角分析\n"
         "      识别转折点与关键决策\n"
         "输出：新 EPISODIC + SEMANTIC\n"
         "      与 SYS2 合并入 BUFFER"),
        ("BUFFER",
         "结构：list of dict\n"
         "类型：episodic / semantic\n"
         "来源：SYS2 (逐事件) + SYS3 (全局)\n"
         "等待：AUDIT 评分后决定去留\n"
         "去向：得分=1 → DB merge"),
        ("AUDIT",
         "输入：BUFFER 候选 + DB 已有经验\n"
         "评分：+1(验证) / -1(证伪) / 0(未测)\n"
         "规则：得分=1 → similarity check\n"
         "      得分≤0 → 丢弃\n"
         "输出：更新 DB (memory.json)"),
    ]
    box_w = Inches(2.05)
    gap = Inches(0.12)
    total_w = len(pipeline) * box_w + (len(pipeline) - 1) * gap
    start_x = (SW - total_w) / 2
    for i, item in enumerate(pipeline):
        title, desc = item[0], item[1]
        x = start_x + i * (box_w + gap)
        add_card_dark(s, x, Inches(1.0), box_w, Inches(3.0))
        add_tb(s, x + Inches(0.08), Inches(1.1), box_w - Inches(0.16), Inches(0.35), title, 13, ORANGE, bold=True)
        add_tb(s, x + Inches(0.08), Inches(1.55), box_w - Inches(0.16), Inches(2.2), desc, 11, WHITE)
        if i < len(pipeline) - 1:
            add_tb(s, x + box_w, Inches(2.3), gap, Inches(0.2), "→", 15, MGRAY, align=PP_ALIGN.CENTER)

    # Bottom: AUDIT (left) + Reflexion (right)
    add_card_dark(s, Inches(0.3), Inches(4.3), Inches(6.0), Inches(2.0))
    add_tb(s, Inches(0.5), Inches(4.4), Inches(5.6), Inches(0.22), "AUDIT 评分细则", 14, ORANGE, bold=True)
    add_tb(s, Inches(0.5), Inches(4.7), Inches(5.6), Inches(1.4),
           "Score = +1 → 该局验证了规则 → 合并到 DB (支持数+1)\n"
           "Score = -1 → 该局证伪了规则 → 丢弃 (矛盾数+1)\n"
           "Score =  0 → 该局未测试 → 丢弃\n"
           "去重引擎: 3层结构化字段匹配 → 相似规则合并计分", 12, WHITE)

    add_card_dark(s, Inches(6.6), Inches(4.3), Inches(6.4), Inches(2.0))
    add_tb(s, Inches(6.8), Inches(4.4), Inches(6.0), Inches(0.22), "借鉴 Reflexion 范式", 14, ORANGE, bold=True)
    add_tb(s, Inches(6.8), Inches(4.7), Inches(6.0), Inches(1.4),
           "Reflexion (Shinn et al., 2023) 提出语言反馈机制：\n"
           "Agent 口头反思任务反馈信号，存入 episodic memory buffer\n"
           "以改进后续决策质量，无需更新模型权重\n\n"
           "我们的赛后反思流水线沿用了这一范式，并引入结构化评分\n"
           "与去重引擎，实现跨对局经验的持续积累与收敛", 12, WHITE)

# ══════════ S7: Memory ══════════
def slide_memory():
    s = make_slide(BACK_DEFAULT, "三·实现 — 三层记忆系统", 7, 'black')
    add_card_dark(s, Inches(0.4), Inches(0.9), Inches(2.5), Inches(0.4))
    add_tb(s, Inches(0.5), Inches(0.93), Inches(2.3), Inches(0.35), "三层结构", 15, ORANGE, bold=True)
    rx = Inches(4.8)
    add_card_dark(s, rx, Inches(0.9), Inches(2.5), Inches(0.4))
    add_tb(s, rx + Inches(0.1), Inches(0.93), Inches(2.3), Inches(0.35), "核心算法", 15, ORANGE, bold=True)
    tiers = [
        ("Working Memory", "最近 100 帧 · TRENDS + DETAIL + DELTA · 单局内, 不持久化"),
        ("Episodic Memory", "跨局 Case 库 (Context + Lesson) + 支持/矛盾计数 · 跨对局 → memory.json"),
        ("Semantic Memory", "去上下文规则库 + 支持/矛盾 + 来源对局 · 跨对局 → memory.json"),
        ("HUMANTIC", "人类先验知识, 硬编码 match-up 指导 · 不参与评分, 仅参考"),
    ]
    for i, (title, desc) in enumerate(tiers):
        y = Inches(1.5) + i * Inches(1.0)
        add_card_dark(s, Inches(0.4), y, Inches(4.0), Inches(0.85))
        add_tb(s, Inches(0.6), y + Inches(0.05), Inches(3.6), Inches(0.22), title, 14, ORANGE, bold=True)
        add_tb(s, Inches(0.6), y + Inches(0.32), Inches(3.6), Inches(0.48), desc, 13, WHITE)
    items = [
        ("去重引擎: 三层结构化字段匹配",
         "Level 1 — 归一化精确匹配 (同义词→术语归一化)\n"
         "Level 2 — 条件前缀匹配 (condition 完全一致)\n"
         "Level 3 — 字段匹配 (skill/action/域/op/val)\n加权评分 ≥0.7 视为等价, 合并计数"),
        ("检索评分: retrieval_score = importance × recency",
         "importance: kill=5, death=4, tower_fall=4\npower_spike=3, gold_spike=2, minion_wave=1\n"
         "recency: 0.5^(days_old / 7) (7天半衰期)\n"
         "按 retrieval_score 降序, top-5 EPISODIC + top-10 SEMANTIC"),
        ("收敛机制 + 跨对局通用记忆",
         "agreement_ratio = supported / (supported + contradicted)\n"
         "低支持率自然沉底, 高支持率优先检索\n"
         "检索时也返回 hero_ai=None 的通用记忆 → 跨英雄学习"),
    ]
    for i, (title, desc) in enumerate(items):
        y = Inches(1.5) + i * Inches(1.5)
        add_card_dark(s, rx, y, Inches(8.0), Inches(1.4))
        add_tb(s, rx + Inches(0.15), y + Inches(0.08), Inches(7.7), Inches(0.22), title, 15, ORANGE, bold=True)
        add_tb(s, rx + Inches(0.15), y + Inches(0.36), Inches(7.7), Inches(0.95), desc, 13, WHITE)

# ══════════ S8: Prompt ══════════
def slide_prompt():
    s = make_slide(BACK_DEFAULT, "三·实现 — Prompt 继承架构与数据模型", 8, 'black')
    add_img(s, PROMPT_PNG, (SW - Inches(6.99)) / 2, Inches(0.7), w=Inches(6.99))

# ══════════ S9: Skill ══════════
def slide_skill():
    s = make_slide(BACK_DEFAULT, "三·实现 — 宏观技能系统 (Skill System)", 9, 'black')
    add_card_dark(s, Inches(0.3), Inches(1.0), Inches(6.0), Inches(2.8))
    add_tb(s, Inches(0.5), Inches(1.1), Inches(5.6), Inches(0.22),
           "Macro Skills (doc → skilldoc → LLM)", 14, ORANGE, bold=True)
    add_tb(s, Inches(0.5), Inches(1.5), Inches(5.6), Inches(2.0),
           "3 个宏观技能, 每个含 3 个子函数:\n"
           "  FARM     last_hit() / move_to_lane() / retreat_to_tower()\n"
           "  POKE     aim_skill() / basic_attack() / reposition_back()\n"
           "  ALL_IN   combo_start() / basic_attack() / chase()\n\n"
           "@register_skill 注册 → skills/__init__.py 自动发现\n"
           "SKILL_REGISTRY.get_doc() 生成 LLM 可见的 skilldoc", 14, WHITE)

    add_card_dark(s, Inches(6.6), Inches(1.0), Inches(6.4), Inches(2.8))
    add_tb(s, Inches(6.8), Inches(1.1), Inches(6.0), Inches(0.22),
           "Concrete Skills (多帧循环执行)", 14, ORANGE, bold=True)
    add_tb(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.0),
           "7 种可执行技能:\n"
           "  FARM(补兵)  POKE(消耗)  ALL_IN(连招)\n"
           "  DEFEND(守塔)  KITE(风筝)  RETREAT(撤退)\n"
           "  PURSUE(追击)  (+MOVE_TO / RECALL 等)\n\n"
           "_start() + update() 循环, 返回 (action, done)\n"
           "ProtocolExecutor 循环调用直到 done=True", 14, WHITE)

    flow = [
        ("LLM", "<action>\n@SKILL_CALL FARM.last_hit()\n</action>"),
        ("ProtocolExecutor\nprocess_batch()", "解析 @SKILL_CALL\n匹配 SKILL_REGISTRY"),
        ("Concrete Skill\n_start() + update()", "循环 update()\n→ (action, done)"),
        ("SkillContext", "make_move() / make_attack()\nmake_skill() / A* 避塔"),
        ("6-tuple", "(btn, mx, mz, sx, sz, target)\nclamp [1, 15]"),
        ("env.step()", "Gamecore 执行\n前进一帧"),
    ]
    bw = Inches(1.9)
    bg = Inches(0.2)
    tw = len(flow) * bw + (len(flow) - 1) * bg
    sx = (SW - tw) / 2
    for i, (title, desc) in enumerate(flow):
        x = sx + i * (bw + bg)
        add_card_dark(s, x, Inches(4.3), bw, Inches(2.0))
        add_tb(s, x + Inches(0.08), Inches(4.4), bw - Inches(0.16), Inches(0.35), title, 14, ORANGE, bold=True)
        add_tb(s, x + Inches(0.08), Inches(4.85), bw - Inches(0.16), Inches(1.3), desc, 14, WHITE)
        if i < len(flow) - 1:
            add_tb(s, x + bw, Inches(5.0), bg, Inches(0.2), "→", 15, MGRAY, align=PP_ALIGN.CENTER)

# ══════════ S10: Demo ══════════
def slide_demo():
    s = make_slide(PROOF_BACK, "四·评估 — 运行演示", 10, 'white')
    # Top-left: trajectory screenshot
    add_img(s, TRAJ_SCREENSHOT, Inches(0.3), Inches(0.85), w=Inches(6.5))
    # Top-right: video render placeholder
    add_card_st(s, Inches(7.2), Inches(0.85), Inches(5.8), Inches(2.8))
    add_tb(s, Inches(7.4), Inches(1.2), Inches(5.4), Inches(2.0),
           "[ 渲染画面 ]\n\n运行实况演示\n\n"
           "gamecore-server (Windows) 渲染游戏画面\n"
           "Docker 容器内运行 LLM Agent\n"
           "每帧决策：parse_state → LLM → 6-tuple → step", 13, DARK, align=PP_ALIGN.CENTER)
    # Bottom-left: trajectory info
    add_card_st(s, Inches(0.3), Inches(4.3), Inches(6.5), Inches(1.8))
    add_tb(s, Inches(0.5), Inches(4.4), Inches(6.1), Inches(0.22), "轨迹回放系统", 14, NAVY, bold=True)
    tb = add_tb(s, Inches(0.5), Inches(4.7), Inches(6.1), Inches(1.2),
               "每局游戏状态以 JSONL 格式完整记录\n"
               "支持逐帧回溯 LLM 的 <think> 推理过程\n"
               "以及 @SKILL_CALL 执行结果\n", 12, DARK)
    add_tb(s, Inches(0.5), Inches(5.5), Inches(6.1), Inches(0.3),
           "网址: https://demo.tom-thu.cn/trajectories/", 12, ORANGE)
    # Bottom-right: stats
    add_card_st(s, Inches(7.2), Inches(4.3), Inches(5.8), Inches(1.8))
    add_tb(s, Inches(7.4), Inches(4.4), Inches(5.4), Inches(0.22), "运行统计", 14, NAVY, bold=True)
    add_tb(s, Inches(7.4), Inches(4.4), Inches(5.4), Inches(1.2),
           "已积累 54 局 JSONL 轨迹日志\n"
           "memory.json 持续跨对局学习\n"
           "支持英雄组合: 后羿/公孙离/马可波罗\n"
           "每次赛后触发 SYS2+SYS3+AUDIT 流水线", 12, DARK)

# ══════════ S11: Comparison ══════════
def slide_comparison():
    s = make_slide(BACK_DEFAULT, "四·评估 — 相关工作对比", 11, 'black')
    headers = ["维度", "TiG (2025)", "WiA-LLM (2026)", "Reflexion (2023)", "Ours"]
    rows = [
        ["方法",  "SFT + GRPO",  "SFT + GRPO",  "语言反馈 + Episodic Buffer", "纯 Prompt 工程"],
        ["训练",  "8×H100 GPU",  "8×H100 GPU",  "零训练",                       "零训练"],
        ["记忆",  "无显式记忆",  "无显式记忆",  "单层 Episodic Buffer",          "三层 + 去重 + 评分"],
        ["决策",  "宏观动作选择", "WhatIf 预测", "反思 → 改进决策",               "WhatIf + @SKILL_CALL"],
        ["可解释","<think> 推理",  "预测文本",    "反思文本",                     "<think> 5段推理"],
    ]
    data = [headers] + rows
    tbl_shape = add_table(s, data, Inches(0.5), Inches(1.1), Inches(12.3), Inches(3.2))
    tbl = tbl_shape.table
    for ri in range(len(data)):
        for ci in range(5):
            c = tbl.cell(ri, ci)
            if ri == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = NAVY if ci < 4 else ORANGE
                tcPr = c._tc.get_or_add_tcPr()
                tcPr.set('anchor', 'ctr')
            else:
                c.fill.solid()
                if ci == 4:
                    c.fill.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xE0)
                elif ri % 2 == 1:
                    c.fill.fore_color.rgb = RGBColor(0xF0, 0xEE, 0xE8)
                else:
                    c.fill.fore_color.rgb = WHITE
                tcPr = c._tc.get_or_add_tcPr()
                tcPr.set('anchor', 'ctr')

    add_card_dark(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.6))
    add_tb(s, Inches(0.7), Inches(4.75), Inches(11.9), Inches(0.25), "核心差异对比", 15, ORANGE, bold=True)
    add_tb(s, Inches(0.7), Inches(5.05), Inches(11.9), Inches(1.2),
           "训练成本:  TiG/WiA-LLM 需 8×H100 GPU + GRPO 训练数千步;  Reflexion 零训练但决策质量受限于基座模型\n"
           "            Ours 零训练, 通过 Prompt 工程 + 记忆系统弥补\n"
           "记忆管理:  TiG/WiA-LLM 无显式记忆;  Reflexion 仅单层 Episodic Buffer;  Ours 三层 + 去重引擎 + 评分\n"
           "决策机制:  TiG 宏观动作分类 / WiA-LLM WhatIf 预测 / Reflexion 反思改进 / Ours WhatIf + @SKILL_CALL", 12, WHITE)

# ══════════ S12: Limitations & Future ══════════
def slide_future():
    s = make_slide(BACK_DEFAULT, "四·评估 — 局限与展望", 12, 'black')
    add_card_dark(s, Inches(0.3), Inches(1.0), Inches(6.0), Inches(4.8))
    add_tb(s, Inches(0.5), Inches(1.1), Inches(5.6), Inches(0.3), "局限", 18, ORANGE, bold=True)
    limitations = [
        ("API 调用成本", "每帧一次 LLM 调用, 长对局 (3200帧) 成本约 $3-5/局, 大规模实验受限。"),
        ("未大规模验证", "当前仅验证射手英雄, 缺乏系统性 Benchmark 测试与胜率统计。"),
        ("单英雄验证", "后羿/公孙离可用, 尚未验证法师、坦克、刺客等其他职业。"),
    ]
    for i, (title, desc) in enumerate(limitations):
        y = Inches(1.65) + i * Inches(0.95)
        add_tb(s, Inches(0.5), y, Inches(5.6), Inches(0.22), f"• {title}", 15, WHITE, bold=True)
        add_tb(s, Inches(0.7), y + Inches(0.28), Inches(5.4), Inches(0.6), desc, 14, RGBColor(0xBB,0xBB,0xBB))

    add_card_dark(s, Inches(6.6), Inches(1.0), Inches(6.4), Inches(4.8))
    add_tb(s, Inches(6.8), Inches(1.1), Inches(6.0), Inches(0.3), "未来方向", 18, ORANGE, bold=True)
    futures = [
        ("多 Agent 合作 (3v3 / 5v5)",
         "扩展到多人对战场景, 多个 LLM Agent 协作。设计通信协议与角色分工。"),
        ("聊天室信息对抗",
         "Agent 不仅在战场对战, 还在聊天室用信息干扰对方。探测意图、制造误导、心理博弈。"),
        ("自动提示调优",
         "引入 BPO / PROMST 等自动优化方法。数据驱动, 无需人工迭代 Prompt。"),
    ]
    for i, (title, desc) in enumerate(futures):
        y = Inches(1.65) + i * Inches(0.95)
        add_tb(s, Inches(6.8), y, Inches(6.0), Inches(0.22), f"• {title}", 15, WHITE, bold=True)
        add_tb(s, Inches(7.0), y + Inches(0.28), Inches(5.8), Inches(0.6), desc, 14, RGBColor(0xBB,0xBB,0xBB))

# ══════════ S13: References ══════════
def slide_refs():
    s = make_slide(LR_BACK, "参考文献", 13, 'white')
    add_card_st(s, Inches(0.3), Inches(0.65), Inches(12.7), Inches(5.8))
    refs = [
        "[1] Hu, S., Huang, T., Liu, G., et al. (2026). A survey on large language model-based game agents. arXiv:2404.02039.",
        "[2] Liao, Y., Gu, Y., Sui, Y., et al. (2025). Think in games via reinforcement learning with large language models. arXiv:2508.21365.",
        "[3] Sui, Y., Zhang, Y., Liao, Y., et al. (2026). What-if analysis of LLMs: Proactive thinking in game worlds. arXiv:2509.04791.",
        "[4] Shinn, N., Cassano, F., Gopinath, A., et al. (2023). Reflexion: Language agents with verbal reinforcement learning. arXiv:2303.11366.",
        "[5] Cheng, J., Liu, X., et al. (2024). Black-box prompt optimization for LLMs. ACL 2024.",
        "[6] Chen, Y., et al. (2024). PROMST: Prompt optimization in multi-step tasks. EMNLP 2024 (Oral).",
        "[7] Pan, R., Xing, S., Diao, S., et al. (2024). Plum: Prompt learning using metaheuristic. ACL 2024 Findings.",
        "[8] Zhan, H., Chen, C., Ding, T., et al. (2024). Zeroth-order black-box prompt tuning. EMNLP 2024 Findings.",
        "[9] Wang, Y., Takanobu, R., Liang, Z., et al. (2025). MEM-alpha: Learning memory construction via RL. arXiv:2509.25911.",
        "[10] Yan, S., Yang, X., Huang, Z., et al. (2025). Memory-R1: LLM agents manage memories via RL. arXiv:2508.19828.",
        "[11] Hu, Y., Liu, S., Yue, Y., et al. (2025). Memory in the age of AI agents: A survey. arXiv:2512.13564.",
        "[12] Du, P. (2026). Memory for autonomous LLM agents: Mechanisms and frontiers. arXiv:2603.07670.",
    ]
    sz, lh = 13, Inches(0.44)
    for i, ref in enumerate(refs):
        add_tb(s, Inches(0.5), Inches(0.85) + i * lh, Inches(12.3), lh, ref, sz, DARK)

# ══════════ S14: Thanks ══════════
def slide_thanks():
    s = make_slide(BACK_DEFAULT, "", 14, 'black')
    add_tb(s, Inches(1.0), Inches(1.0), Inches(11.3), Inches(0.8),
           "感谢聆听", 40, WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_tb(s, Inches(1.0), Inches(1.9), Inches(11.3), Inches(0.5),
           "Honor of Kings LLM Agent", 18, ORANGE, align=PP_ALIGN.CENTER)

    add_card_dark(s, Inches(1.5), Inches(2.8), Inches(4.5), Inches(3.0))
    add_tb(s, Inches(1.7), Inches(3.0), Inches(4.1), Inches(0.25), "杨坤鑫", 14, ORANGE, bold=True)
    add_tb(s, Inches(1.7), Inches(3.4), Inches(4.1), Inches(2.2),
           "- Pipeline 设计与搭建\n- Prompt 体系设计\n- 记忆系统架构与实现\n- Trajectory 网页部署", 14, WHITE)

    add_card_dark(s, Inches(7.3), Inches(2.8), Inches(4.5), Inches(3.0))
    add_tb(s, Inches(7.5), Inches(3.0), Inches(4.1), Inches(0.25), "张腾达", 14, ORANGE, bold=True)
    add_tb(s, Inches(7.5), Inches(3.4), Inches(4.1), Inches(2.2),
           "- 宏观技能迭代 (FARM/POKE等)\n- 记忆系统去重引擎\n- Pipeline 调试与对接\n- 记忆系统调试与测试", 14, WHITE)

    add_tb(s, Inches(3.0), Inches(6.3), Inches(7.3), Inches(0.3),
           "欢迎批评指正", 14, MGRAY, align=PP_ALIGN.CENTER)

# ══════════ Build ══════════
print("Building 14 slides...", flush=True)
slides = [
    (slide_title, "Title"),
    (slide_toc, "TOC"),
    (slide_problem, "Problem & Related Work"),
    (slide_arch, "Architecture"),
    (slide_sys1, "SYS1 Decision"),
    (slide_reflection, "Reflection Pipeline"),
    (slide_memory, "Memory System"),
    (slide_prompt, "Prompt Architecture"),
    (slide_skill, "Skill System"),
    (slide_demo, "Demo"),
    (slide_comparison, "Comparison"),
    (slide_future, "Limitations & Future"),
    (slide_refs, "References"),
    (slide_thanks, "Thanks"),
]
for i, (fn, name) in enumerate(slides):
    fn()
    print(f"  {i+1:2d}/14: {name}", flush=True)

prs.save(OUTPUT)
print(f"\nSaved: {OUTPUT}", flush=True)
print(json.dumps({"status": "ok", "path": OUTPUT}))
